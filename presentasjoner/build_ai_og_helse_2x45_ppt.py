from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUTPUT_PATH = Path(__file__).resolve().parent / "AI_og_Helse_2x45min_klinisk_forelesning.pptx"


PURPLE = RGBColor(125, 60, 152)
PURPLE_DARK = RGBColor(91, 44, 111)
BLUE = RGBColor(46, 134, 193)
TEXT = RGBColor(31, 26, 23)
MUTED = RGBColor(93, 82, 75)
BG = RGBColor(251, 248, 239)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(226, 216, 206)


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.32)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = PURPLE
    band.line.fill.background()


def add_title(slide, title, subtitle=None, kicker=None):
    add_bg(slide)

    if kicker:
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(6.0), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = kicker
        r.font.name = "Aptos"
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = PURPLE_DARK

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(11.7), Inches(1.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = TEXT

    if subtitle:
        stb = slide.shapes.add_textbox(Inches(0.75), Inches(2.1), Inches(10.8), Inches(1.2))
        stf = stb.text_frame
        stf.word_wrap = True
        p = stf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = "Aptos"
        r.font.size = Pt(18)
        r.font.color.rgb = MUTED


def add_footer_reference(slide, references):
    if not references:
        return

    box = slide.shapes.add_textbox(Inches(0.75), Inches(6.95), Inches(11.8), Inches(0.4))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Referanser i repoet: " + " | ".join(references)
    r.font.name = "Aptos"
    r.font.size = Pt(9)
    r.font.italic = True
    r.font.color.rgb = MUTED


def add_content_slide(slide, title, bullets, references=None, section=None):
    add_bg(slide)

    if section:
        sec = slide.shapes.add_textbox(Inches(0.75), Inches(0.55), Inches(4.5), Inches(0.35))
        p = sec.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = section
        r.font.name = "Aptos"
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = PURPLE_DARK

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.95), Inches(11.6), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = TEXT

    content_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.85), Inches(11.8), Inches(4.8)
    )
    content_shape.fill.solid()
    content_shape.fill.fore_color.rgb = WHITE
    content_shape.line.color.rgb = LINE

    text_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.08), Inches(11.1), Inches(4.3))
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP

    first = True
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(6)
        p.bullet = True
        r = p.add_run()
        r.text = bullet
        r.font.name = "Aptos"
        r.font.size = Pt(20)
        r.font.color.rgb = TEXT
        first = False

    add_footer_reference(slide, references)


def add_section_slide(slide, title, subtitle):
    add_bg(slide)

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(1.2), Inches(11.4), Inches(4.8)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = LINE

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.95), Inches(1.2), Inches(0.22), Inches(4.8)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = PURPLE
    accent.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(9.8), Inches(1.0))
    p = tbox.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = TEXT

    sbox = slide.shapes.add_textbox(Inches(1.5), Inches(3.0), Inches(9.4), Inches(1.4))
    p = sbox.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = subtitle
    r.font.name = "Aptos"
    r.font.size = Pt(20)
    r.font.color.rgb = MUTED


def add_repo_overview_slide(slide, title, items, references=None, section=None):
    add_bg(slide)

    if section:
        sec = slide.shapes.add_textbox(Inches(0.75), Inches(0.55), Inches(4.5), Inches(0.35))
        p = sec.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = section
        r.font.name = "Aptos"
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = PURPLE_DARK

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.95), Inches(11.6), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = TEXT

    x_positions = [0.75, 4.55, 8.35]
    y_positions = [1.9, 4.2]
    width = 3.45
    height = 1.95

    for idx, item in enumerate(items):
        col = idx % 3
        row = idx // 3
        x = Inches(x_positions[col])
        y = Inches(y_positions[row])

        panel = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, Inches(width), Inches(height)
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = WHITE
        panel.line.color.rgb = LINE

        icon = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.18), Inches(0.45), Inches(0.35))
        p = icon.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = item["icon"]
        r.font.name = "Aptos"
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = PURPLE_DARK

        tb = slide.shapes.add_textbox(x + Inches(0.62), y + Inches(0.18), Inches(2.55), Inches(0.35))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = item["title"]
        r.font.name = "Aptos"
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = TEXT

        db = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.62), Inches(3.0), Inches(0.78))
        tf = db.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = item["description"]
        r.font.name = "Aptos"
        r.font.size = Pt(12.5)
        r.font.color.rgb = MUTED

        nb = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(1.42), Inches(3.0), Inches(0.3))
        p = nb.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = item["notebook"]
        r.font.name = "Aptos"
        r.font.size = Pt(10)
        r.font.italic = True
        r.font.color.rgb = BLUE

    add_footer_reference(slide, references)


def add_case_slide(slide, title, scenario, left_title, left_points, right_title, right_points, references=None, section=None):
    add_bg(slide)

    if section:
        sec = slide.shapes.add_textbox(Inches(0.75), Inches(0.55), Inches(4.5), Inches(0.35))
        p = sec.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = section
        r.font.name = "Aptos"
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = PURPLE_DARK

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.95), Inches(11.4), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = TEXT

    scenario_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.8), Inches(11.8), Inches(1.1)
    )
    scenario_box.fill.solid()
    scenario_box.fill.fore_color.rgb = WHITE
    scenario_box.line.color.rgb = LINE

    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.03), Inches(11.2), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = scenario
    r.font.name = "Aptos"
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = TEXT

    for x, panel_title, panel_points, symbol in [
        (0.75, left_title, left_points, "⊕"),
        (6.75, right_title, right_points, "◇"),
    ]:
        panel = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.15), Inches(5.8), Inches(3.25)
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = WHITE
        panel.line.color.rgb = LINE

        sb = slide.shapes.add_textbox(Inches(x + 0.25), Inches(3.38), Inches(0.4), Inches(0.3))
        p = sb.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = symbol
        r.font.name = "Aptos"
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = PURPLE_DARK

        hb = slide.shapes.add_textbox(Inches(x + 0.65), Inches(3.35), Inches(4.6), Inches(0.35))
        p = hb.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = panel_title
        r.font.name = "Aptos"
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = TEXT

        pb = slide.shapes.add_textbox(Inches(x + 0.28), Inches(3.82), Inches(5.1), Inches(2.2))
        tf = pb.text_frame
        tf.word_wrap = True
        first = True
        for point in panel_points:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.bullet = True
            p.level = 0
            p.space_after = Pt(4)
            r = p.add_run()
            r.text = point
            r.font.name = "Aptos"
            r.font.size = Pt(15)
            r.font.color.rgb = TEXT
            first = False

    add_footer_reference(slide, references)


def add_title_slide(slide, title, subtitle, footer):
    add_bg(slide)

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.95), Inches(0.85), Inches(4.45), Inches(5.5)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = WHITE
    accent.line.color.rgb = LINE

    deco = slide.shapes.add_textbox(Inches(8.35), Inches(1.2), Inches(3.6), Inches(0.6))
    p = deco.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = "◎   ◇   ↗"
    r.font.name = "Aptos"
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = PURPLE_DARK

    tbox = slide.shapes.add_textbox(Inches(0.85), Inches(1.25), Inches(6.6), Inches(2.1))
    tf = tbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = TEXT

    sbox = slide.shapes.add_textbox(Inches(0.9), Inches(3.35), Inches(6.2), Inches(1.3))
    p = sbox.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = subtitle
    r.font.name = "Aptos"
    r.font.size = Pt(20)
    r.font.color.rgb = MUTED

    fbox = slide.shapes.add_textbox(Inches(0.9), Inches(6.55), Inches(8.0), Inches(0.45))
    p = fbox.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = footer
    r.font.name = "Aptos"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = PURPLE_DARK

    info_box = slide.shapes.add_textbox(Inches(8.35), Inches(2.0), Inches(3.2), Inches(3.3))
    tf = info_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(
        [
            "Målgruppe",
            "Leger, sykepleiere, psykologer, fysioterapeuter, radiografer, bioingeniører",
            "",
            "Format",
            "2 x 45 minutter",
            "",
            "Utgangspunkt",
            "Kurset AI og Helse",
        ]
    ):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.name = "Aptos"
        r.font.size = Pt(15 if line and line[0].isupper() and len(line.split()) <= 2 else 13)
        r.font.bold = line in {"Målgruppe", "Format", "Utgangspunkt"}
        r.font.color.rgb = TEXT if r.font.bold else MUTED


SLIDES = [
    {
        "type": "title",
        "title": "AI i helse: hva er nyttig, hva er hype, og hva må vi forstå?",
        "subtitle": "To forelesningsøkter à 45 minutter for klinisk målgruppe med liten eller moderat teknisk bakgrunn.",
        "footer": "Basert på AI og Helse-repoet og tilhørende notebooks",
    },
    {
        "type": "content",
        "section": "Innledning",
        "title": "Hvem denne forelesningen er for",
        "bullets": [
            "Leger, sykepleiere, psykologer, fysioterapeuter, radiografer og bioingeniører",
            "Andre som jobber klinisk eller tett på helsefaglige arbeidsprosesser",
            "Personer som vil forstå AI uten å måtte gå dypt inn i matematikk eller programmering",
        ],
        "references": ["README.md", "docs/om-kurset.html"],
    },
    {
        "type": "content",
        "section": "Innledning",
        "title": "Hva dere skal sitte igjen med",
        "bullets": [
            "Et språk for å snakke om AI i helse",
            "En forståelse av hva AI er god og mindre god til",
            "Noen konkrete spørsmål å stille før man tar AI i bruk",
            "En bedre forståelse av risiko, ansvar og klinisk nytte",
        ],
    },
    {
        "type": "content",
        "section": "Innledning",
        "title": "Hvorfor AI i helse er et viktig tema nå",
        "bullets": [
            "Data, bilder, signaler og tekst finnes i stor skala",
            "Arbeidshverdagen er presset og kompleks",
            "Mange AI-verktøy lover effektivitet og bedre beslutninger",
            "Men helsefeltet tåler dårlig feil og falsk trygghet",
        ],
    },
    {
        "type": "content",
        "section": "Innledning",
        "title": "Hvorfor AI i helse er vanskeligere enn i mange andre felt",
        "bullets": [
            "Pasientsikkerhet og ansvar står sentralt",
            "Kontekst betyr mye mer enn i mange generelle AI-demoer",
            "Skjevheter og generalisering er klinisk viktige spørsmål",
            "Arbeidsflyt og implementering avgjør om noe faktisk fungerer",
        ],
    },
    {
        "type": "content",
        "section": "Innledning",
        "title": "Et raskt startspørsmål til salen",
        "bullets": [
            "Hvem har brukt ChatGPT eller lignende verktøy?",
            "Hvem har møtt AI i journalsystemer, bildeanalyse eller beslutningsstøtte?",
            "Hvem er nysgjerrig, men usikker på hva som faktisk er nyttig?",
        ],
    },
    {
        "type": "content",
        "section": "Del 1",
        "title": "Hva mener vi egentlig med AI?",
        "bullets": [
            "AI er et samlebegrep, ikke én enkelt teknologi",
            "AI kan brukes til å klassifisere, predikere, anbefale eller generere",
            "I helse er det viktig å skille mellom teknologitype og faktisk bruksområde",
        ],
        "references": ["uke01-introduksjon/02-hva-er-ai.ipynb"],
    },
    {
        "type": "content",
        "section": "Del 1",
        "title": "Fra regler til læring",
        "bullets": [
            "Tidlige systemer bygget på eksplisitte regler",
            "Moderne systemer lærer mønstre fra data",
            "Dette gir mer fleksibilitet, men også mindre gjennomsiktighet",
        ],
        "references": ["uke01-introduksjon/05-regelbaserte-systemer.ipynb"],
    },
    {
        "type": "content",
        "section": "Del 1",
        "title": "Tre hovedfamilier vi møter i kurset",
        "bullets": [
            "Klassisk maskinlæring",
            "Dyplæring",
            "Generativ AI",
        ],
        "references": ["uke01-introduksjon/04-ai-ml-dl-forskjeller.ipynb"],
    },
    {
        "type": "content",
        "section": "Del 1",
        "title": "Klassisk maskinlæring, veldig kort",
        "bullets": [
            "Modeller finner mønstre i strukturerte data",
            "Brukes ofte til prediksjon og klassifikasjon",
            "Relevante spørsmål er hvilke data som brukes, hva som måles og hvilke feil som betyr mest",
        ],
        "references": ["uke02-klassisk-ml/01-klassisk-ml-101.ipynb"],
    },
    {
        "type": "content",
        "section": "Del 1",
        "title": "Dyplæring, veldig kort",
        "bullets": [
            "Særlig nyttig når data er komplekse",
            "Typiske eksempler er bilder, signaler og tekst",
            "Mer kraftfullt, men ofte vanskeligere å forklare og kontrollere",
        ],
        "references": [
            "uke03-dyplæring/01a_nn_intro.ipynb",
            "uke03-dyplæring/02a_cnn_bildeklassifikasjon.ipynb",
        ],
    },
    {
        "type": "content",
        "section": "Del 1",
        "title": "Generativ AI, veldig kort",
        "bullets": [
            "Lager tekst, oppsummeringer, forslag eller multimodalt innhold",
            "Kan være nyttig i dokumentasjon, dialog og kunnskapsstøtte",
            "Men kan også hallusinere og virke mer sikker enn den burde",
        ],
        "references": [
            "uke04-generativ-ai/02_llm_grunnleggende.ipynb",
            "uke04-generativ-ai/03_prompt_engineering.ipynb",
        ],
    },
    {
        "type": "content",
        "section": "Del 1",
        "title": "Hvorfor helse er interessant for AI",
        "bullets": [
            "Mange mønstre er vanskelige å se uten støtte",
            "Kliniske beslutninger er ofte probabilistiske",
            "Bilder, signaler og tekst spiller en stor rolle",
            "Små forbedringer kan ha stor verdi i praksis",
        ],
    },
    {
        "type": "content",
        "section": "Del 1",
        "title": "Hvorfor helse også er risikofylt for AI",
        "bullets": [
            "Kontekst betyr mye",
            "Data er ofte skjeve, mangelfulle eller lokale",
            "Feil kan skade pasienter",
            "Tillit og ansvar kan ikke automatiseres bort",
        ],
    },
    {
        "type": "content",
        "section": "Del 1",
        "title": "Eksempel 1: bildeanalyse",
        "bullets": [
            "Radiologi, patologi og dermatologi bygger alle på mønstergjenkjenning",
            "AI kan være sterk på å oppdage signaler i bilder",
            "Det avgjørende spørsmålet er hvordan dette brukes i klinisk arbeid",
        ],
        "references": ["uke03-dyplæring/03_medisinsk_bildeklassifikasjon_MR.ipynb"],
    },
    {
        "type": "content",
        "section": "Del 1",
        "title": "Eksempel 2: risikomodeller",
        "bullets": [
            "Hvem har økt risiko?",
            "Hvem bør følges opp tettere?",
            "Når blir en prediksjon faktisk nyttig for klinisk beslutningsstøtte?",
        ],
        "references": ["uke06-klinisk-praksis/01_risikomodell_logistisk_regresjon_kalibrering_shap.ipynb"],
    },
    {
        "type": "content",
        "section": "Del 1",
        "title": "Eksempel 3: språk og dokumentasjon",
        "bullets": [
            "Oppsummering av tekst og strukturering av informasjon",
            "Forslag til formuleringer og støtte i administrativt arbeid",
            "Lavere terskel enn mange andre AI-bruksområder, men fortsatt ikke risikofritt",
        ],
        "references": [
            "uke04-generativ-ai/04_chatgpt_claude_api.ipynb",
            "intro_openai_anthropic.ipynb",
        ],
    },
    {
        "type": "content",
        "section": "Del 1",
        "title": "Eksempel 4: velferdsteknologi og sensorer",
        "bullets": [
            "Tolkning av aktivitet og avvik",
            "Varsling og prioritering i hjem og omsorg",
            "Teknologi som virker tett på hverdagsliv og omsorgspraksis",
        ],
        "references": [
            "uke07-velferdsteknologi/02_sensorer_aktivitet_og_hendelsesforståelse.ipynb",
            "uke07-velferdsteknologi/03_beslutningsstøtte_i_hjem_og_omsorg.ipynb",
        ],
    },
    {
        "type": "content",
        "section": "Del 1",
        "title": "Eksempel 5: generativ AI som støtte",
        "bullets": [
            "Utkast til tekst og oppsummering av informasjon",
            "Strukturering av faglig innhold og spørsmål",
            "Kan være nyttig som støtte, men ikke som automatisk fasit",
        ],
    },
    {
        "type": "repo_overview",
        "section": "Del 1",
        "title": "Temaene i repoet, som kurskart",
        "items": [
            {
                "icon": "◎",
                "title": "Uke 1",
                "description": "Begreper, historikk og hva AI faktisk er i helsefaglig sammenheng.",
                "notebook": "uke01-introduksjon/02-hva-er-ai.ipynb",
            },
            {
                "icon": "∿",
                "title": "Uke 2",
                "description": "Klassisk maskinlæring, prediksjon og grunnleggende evaluering.",
                "notebook": "uke02-klassisk-ml/01-klassisk-ml-101.ipynb",
            },
            {
                "icon": "▣",
                "title": "Uke 3",
                "description": "Dyplæring, bilder, signaler og forklarbarhet.",
                "notebook": "uke03-dyplæring/03_medisinsk_bildeklassifikasjon_MR.ipynb",
            },
            {
                "icon": "✦",
                "title": "Uke 4-5",
                "description": "Generativ AI, språkmodeller, prompt engineering og agentisk AI.",
                "notebook": "uke04-generativ-ai/02_llm_grunnleggende.ipynb",
            },
            {
                "icon": "⊕",
                "title": "Uke 6",
                "description": "Klinisk praksis, risikomodeller og beslutningsstøtte.",
                "notebook": "uke06-klinisk-praksis/01_risikomodell_logistisk_regresjon_kalibrering_shap.ipynb",
            },
            {
                "icon": "◇",
                "title": "Uke 7-8",
                "description": "Velferdsteknologi, etikk, regulering og trustworthy AI.",
                "notebook": "uke08-etikk-implementering/05_trustworthy_ai_i_helse.ipynb",
            },
        ],
        "references": ["README.md"],
    },
    {
        "type": "case",
        "section": "Del 1",
        "title": "Visuelt case: AI i radiologi",
        "scenario": "Tenk deg en radiologisk arbeidsflyt der et AI-system markerer mistenkelige funn før klinikeren leser bildet.",
        "left_title": "Mulig verdi",
        "left_points": [
            "Kan bidra til raskere sortering og prioritering",
            "Kan være ekstra nyttig i store volumer og ved belastet arbeidsflyt",
            "Kan fungere som et ekstra blikk, ikke nødvendigvis en erstatning",
        ],
        "right_title": "Viktige spørsmål",
        "right_points": [
            "Hvordan er systemet validert og i hvilken setting?",
            "Hvilke typer feil gjør det oftest?",
            "Hvordan påvirker det klinikerens oppmerksomhet og beslutning?",
        ],
        "references": ["uke03-dyplæring/03_medisinsk_bildeklassifikasjon_MR.ipynb"],
    },
    {
        "type": "content",
        "section": "Del 1",
        "title": "Hva AI ofte er god på",
        "bullets": [
            "Mønstergjenkjenning i store datamengder",
            "Skalerbarhet og repetitive oppgaver",
            "Forslag, sortering og raske førsteutkast",
        ],
    },
    {
        "type": "content",
        "section": "Del 1",
        "title": "Hva AI ofte er dårlig på",
        "bullets": [
            "Dyp klinisk kontekstforståelse",
            "Å vite hva som ikke er representert i data",
            "Å forstå ansvar, konsekvenser og etiske avveininger",
            "Å være pålitelig i nye situasjoner",
        ],
    },
    {
        "type": "content",
        "section": "Del 1",
        "title": "Hva klinikere bør spørre om først",
        "bullets": [
            "Hvilket problem er dette ment å løse?",
            "Hva er input og hva er output?",
            "Hvem påvirkes hvis systemet tar feil?",
            "Hva skjer i praksis hvis vi følger anbefalingen?",
        ],
    },
    {
        "type": "content",
        "section": "Del 1",
        "title": "Mini-case 1: høy risiko for sepsis",
        "bullets": [
            "Et AI-system markerer høy risiko for sepsis",
            "Hva vil du vite om datagrunnlaget?",
            "Hva vil du vite om feilmarginer og terskler?",
            "Hvordan bør dette brukes i klinisk arbeidsflyt?",
        ],
    },
    {
        "type": "content",
        "section": "Del 1",
        "title": "Fra fascinasjon til vurdering",
        "bullets": [
            "Ikke spør bare om systemet fungerer",
            "Spør for hvem det fungerer, i hvilken setting og med hvilke konsekvenser",
            "Spør hva som er baseline, og hva som faktisk blir bedre",
        ],
    },
    {
        "type": "content",
        "section": "Del 1",
        "title": "Oppsummering av del 1",
        "bullets": [
            "AI er ikke én ting, men flere familier av verktøy og modeller",
            "Helse er et felt med både store muligheter og høy risiko",
            "Klinisk nytte avhenger av kontekst, ikke bare teknologi",
        ],
    },
    {
        "type": "section",
        "title": "Del 2",
        "subtitle": "Fra teknologi til trygg, nyttig og ansvarlig bruk i praksis",
    },
    {
        "type": "content",
        "section": "Del 2",
        "title": "Hva betyr det at en modell virker?",
        "bullets": [
            "Høy ytelse i et datasett er ikke nok",
            "Vi må spørre om klinisk nytte, robusthet og praktisk bruk",
            "God teknologi kan fortsatt være dårlig implementert",
        ],
    },
    {
        "type": "content",
        "section": "Del 2",
        "title": "Nøyaktighet er ikke hele historien",
        "bullets": [
            "Sensitivitet og spesifisitet sier ulike ting",
            "Falske positive og falske negative har ulike kliniske kostnader",
            "Terskler og konsekvenser må forstås før verktøy tas i bruk",
        ],
        "references": ["uke06-klinisk-praksis/02_klinisk_beslutningsstøtte_terskler_og_avveininger.ipynb"],
    },
    {
        "type": "content",
        "section": "Del 2",
        "title": "Kalibrering på klinisk språk",
        "bullets": [
            "Hvis modellen sier 10 prosent risiko, stemmer det omtrent?",
            "Kalibrering er viktig når tall skal brukes i beslutningsstøtte",
            "Dårlig kalibrering kan gi falsk trygghet eller unødig alarm",
        ],
        "references": ["uke06-klinisk-praksis/01_risikomodell_logistisk_regresjon_kalibrering_shap.ipynb"],
    },
    {
        "type": "content",
        "section": "Del 2",
        "title": "Hvorfor generalisering er vanskelig",
        "bullets": [
            "En modell kan fungere godt ett sted og dårlig et annet",
            "Pasientgrupper, utstyr og praksis kan være ulike",
            "Data i drift ligner ikke alltid data i utvikling",
        ],
        "references": ["uke06-klinisk-praksis/03_validering_generalisering_og_subgrupper.ipynb"],
    },
    {
        "type": "content",
        "section": "Del 2",
        "title": "Distribusjonsskifte",
        "bullets": [
            "Pasientpopulasjoner endrer seg",
            "Arbeidsflyt og rutiner endrer seg",
            "Data endrer karakter over tid",
            "Modeller må derfor monitoreres, ikke bare lanseres",
        ],
    },
    {
        "type": "content",
        "section": "Del 2",
        "title": "Subgrupper og rettferdighet",
        "bullets": [
            "Samme modell kan ha ulik ytelse for ulike grupper",
            "Dette er både et faglig og et etisk spørsmål",
            "Lik gjennomsnittsytelse er ikke det samme som rettferdig ytelse",
        ],
        "references": ["uke08-etikk-implementering/02_bias_rettferdighet.ipynb"],
    },
    {
        "type": "content",
        "section": "Del 2",
        "title": "Hvor kan bias oppstå?",
        "bullets": [
            "I hvem som er med i datasettet",
            "I hvordan data måles og merkes",
            "I hvordan modellen utvikles",
            "I hvordan resultatene brukes og tolkes i praksis",
        ],
    },
    {
        "type": "content",
        "section": "Del 2",
        "title": "Forklarbarhet: nyttig, men ikke magisk",
        "bullets": [
            "Forklaringer kan hjelpe oss å se hva modellen reagerer på",
            "Men forklaringer beviser ikke at modellen er klinisk riktig",
            "Forklarbarhet må brukes kritisk, ikke ritualistisk",
        ],
        "references": ["uke03-dyplæring/02c_cnn_testing.ipynb"],
    },
    {
        "type": "content",
        "section": "Del 2",
        "title": "Human-in-the-loop",
        "bullets": [
            "AI bør vanligvis støtte, ikke erstatte",
            "Klinikerens vurdering forsvinner ikke",
            "Det avgjørende er hvordan menneske og system samspiller",
        ],
    },
    {
        "type": "content",
        "section": "Del 2",
        "title": "Arbeidsflyt er en del av kvaliteten",
        "bullets": [
            "En god modell som ikke passer i praksis har liten verdi",
            "Dårlig integrasjon kan skape friksjon og nye feil",
            "Implementering er et klinisk og organisatorisk spørsmål",
        ],
        "references": ["uke06-klinisk-praksis/04_fra_modell_til_klinisk_arbeidsflyt.ipynb"],
    },
    {
        "type": "content",
        "section": "Del 2",
        "title": "Hva skjer hvis modellen er feil?",
        "bullets": [
            "Hvem merker det?",
            "Hvem stopper det?",
            "Hvor raskt oppdages problemet?",
            "Hvem har myndighet til å korrigere eller slå av systemet?",
        ],
    },
    {
        "type": "content",
        "section": "Del 2",
        "title": "Personvern og helseopplysninger",
        "bullets": [
            "Helseopplysninger krever særlig aktsomhet",
            "Dataminimering, tilgangskontroll og informasjonsplikt betyr noe",
            "Tillit er en del av implementeringen",
        ],
        "references": ["uke08-etikk-implementering/01_gdpr_personvern.ipynb"],
    },
    {
        "type": "content",
        "section": "Del 2",
        "title": "Etikk i praksis",
        "bullets": [
            "Autonomi",
            "Velgjørenhet",
            "Ikke-skade",
            "Rettferdighet",
        ],
        "references": ["uke08-etikk-implementering/04_ai_etikk_i_medisinen.ipynb"],
    },
    {
        "type": "content",
        "section": "Del 2",
        "title": "Regulering i grove trekk",
        "bullets": [
            "Noe AI i helse kan være medisinsk utstyr",
            "Da følger dokumentasjon, vurdering og ansvar",
            "Regulering er ikke bare byråkrati, men en del av sikkerheten",
        ],
        "references": ["uke08-etikk-implementering/03_ce_mdr_regulering.ipynb"],
    },
    {
        "type": "content",
        "section": "Del 2",
        "title": "Hva betyr trustworthy AI i helse?",
        "bullets": [
            "Robusthet",
            "Transparens",
            "Menneskelig kontroll",
            "Validering og monitorering over tid",
        ],
        "references": ["uke08-etikk-implementering/05_trustworthy_ai_i_helse.ipynb"],
    },
    {
        "type": "content",
        "section": "Del 2",
        "title": "Generativ AI i klinikken: egne utfordringer",
        "bullets": [
            "Hallusinasjoner og overbevisende språk",
            "Uklare kilder og usikker sporbarhet",
            "Vanskelig å vite når modellen gjetter",
            "Stor risiko for ukritisk overtakelse av formuleringer",
        ],
        "references": [
            "uke04-generativ-ai/02_llm_grunnleggende.ipynb",
            "uke04-generativ-ai/04_chatgpt_claude_api.ipynb",
        ],
    },
    {
        "type": "content",
        "section": "Del 2",
        "title": "Hva bør man ikke outsource blindt?",
        "bullets": [
            "Diagnostiske konklusjoner",
            "Høyrisikobeslutninger",
            "Følsom pasientkommunikasjon",
            "Ansvaret for klinisk skjønn",
        ],
    },
    {
        "type": "repo_overview",
        "section": "Del 2",
        "title": "Hvor i repoet finner dere temaene vi har snakket om?",
        "items": [
            {
                "icon": "⊕",
                "title": "Risiko og kalibrering",
                "description": "Hvordan prediksjon blir til klinisk vurdering og terskelvalg.",
                "notebook": "uke06-klinisk-praksis/02_klinisk_beslutningsstøtte_terskler_og_avveininger.ipynb",
            },
            {
                "icon": "∿",
                "title": "Generalisering",
                "description": "Hvorfor modeller feiler når data og kontekst endrer seg.",
                "notebook": "uke06-klinisk-praksis/03_validering_generalisering_og_subgrupper.ipynb",
            },
            {
                "icon": "◇",
                "title": "Bias og rettferdighet",
                "description": "Hvordan skjevhet kan oppstå og hvorfor det betyr noe klinisk.",
                "notebook": "uke08-etikk-implementering/02_bias_rettferdighet.ipynb",
            },
            {
                "icon": "◎",
                "title": "Personvern",
                "description": "GDPR og praktiske vurderinger ved helseopplysninger.",
                "notebook": "uke08-etikk-implementering/01_gdpr_personvern.ipynb",
            },
            {
                "icon": "↗",
                "title": "Regulering",
                "description": "CE, MDR og ansvar ved innføring av medisinsk AI.",
                "notebook": "uke08-etikk-implementering/03_ce_mdr_regulering.ipynb",
            },
            {
                "icon": "✦",
                "title": "Trustworthy AI",
                "description": "Robusthet, monitorering og trygg innføring i praksis.",
                "notebook": "uke08-etikk-implementering/05_trustworthy_ai_i_helse.ipynb",
            },
        ],
        "references": ["README.md", "uke08-etikk-implementering/README.md"],
    },
    {
        "type": "content",
        "section": "Del 2",
        "title": "Hva kan være lavere risiko og likevel nyttig?",
        "bullets": [
            "Strukturering av tekst",
            "Oppsummering av informasjon",
            "Administrative støtteoppgaver",
            "Utkast som alltid kontrolleres av fagperson",
        ],
    },
    {
        "type": "case",
        "section": "Del 2",
        "title": "Visuelt case: generativ AI i dokumentasjon",
        "scenario": "Tenk deg at et system lager utkast til epikrise, pasientbrev eller journalsammendrag basert på eksisterende tekst og struktur.",
        "left_title": "Mulige gevinster",
        "left_points": [
            "Sparer tid i formulering og strukturering",
            "Kan gjøre dokumentasjon mer konsistent",
            "Kan være nyttig som førsteutkast i lavere risiko-oppgaver",
        ],
        "right_title": "Hva må kontrolleres",
        "right_points": [
            "At innholdet faktisk stemmer med pasientforløpet",
            "At det ikke introduseres feil, hallusinasjoner eller tvetydighet",
            "At kliniker fortsatt har ansvar for språk, innhold og konsekvenser",
        ],
        "references": ["uke04-generativ-ai/04_chatgpt_claude_api.ipynb", "intro_openai_anthropic.ipynb"],
    },
    {
        "type": "content",
        "section": "Del 2",
        "title": "Praktisk sjekkliste før innføring",
        "bullets": [
            "Hvilket problem skal løses?",
            "Finnes det en god baseline uten AI?",
            "Er data gode nok og relevant validering gjort?",
            "Hvem har ansvar hvis noe går galt?",
        ],
    },
    {
        "type": "content",
        "section": "Del 2",
        "title": "Praktisk sjekkliste under bruk",
        "bullets": [
            "Følger vi med på avvik?",
            "Forstår brukerne begrensningene?",
            "Har vi tydelig rollefordeling?",
            "Vet vi når systemet ikke bør brukes?",
        ],
    },
    {
        "type": "content",
        "section": "Del 2",
        "title": "Praktisk sjekkliste etter innføring",
        "bullets": [
            "Re-evaluer ytelse over tid",
            "Se etter endringer i pasientgrunnlag og praksis",
            "Samle tilbakemeldinger fra brukere",
            "Oppdater eller stopp løsninger som ikke fungerer godt nok",
        ],
    },
    {
        "type": "content",
        "section": "Del 2",
        "title": "Mini-case 2: generativ AI lager pasientbrevutkast",
        "bullets": [
            "Hva kan være nyttig i et slikt verktøy?",
            "Hva må alltid kontrolleres manuelt?",
            "Hvor kan feil oppstå?",
            "Hvordan kan dette brukes uten å svekke kvalitet og tillit?",
        ],
    },
    {
        "type": "content",
        "section": "Avslutning",
        "title": "Tre ting klinikere bør ta med seg",
        "bullets": [
            "AI kan være nyttig, men må brukes med forståelse",
            "God ytelse er ikke nok uten klinisk kontekst",
            "Kritiske spørsmål er en del av faglig ansvar",
        ],
    },
    {
        "type": "content",
        "section": "Avslutning",
        "title": "Tre ting organisasjoner bør ta med seg",
        "bullets": [
            "Implementering handler om mer enn teknologi",
            "Arbeidsflyt, opplæring og ansvar er avgjørende",
            "Trygg bruk krever monitorering, ikke bare lansering",
        ],
    },
    {
        "type": "content",
        "section": "Avslutning",
        "title": "Hvordan bruke repoet videre etter forelesningen",
        "bullets": [
            "Start med uke 1 for begreper og oversikt",
            "Bruk uke 4 for generativ AI og uke 6 for klinisk beslutningsstøtte",
            "Bruk uke 8 når dere vil arbeide med etikk, regulering og trustworthy AI",
        ],
        "references": [
            "README.md",
            "uke04-generativ-ai/README.md",
            "uke06-klinisk-praksis/README.md",
            "uke08-etikk-implementering/README.md",
        ],
    },
    {
        "type": "content",
        "section": "Avslutning",
        "title": "Spørsmål til diskusjon i salen",
        "bullets": [
            "Hvor ser dere størst nytte i egen praksis?",
            "Hvor ser dere størst risiko?",
            "Hva ville dere krevd å få vite før dere stolte på et AI-system?",
        ],
    },
    {
        "type": "content",
        "section": "Avslutning",
        "title": "Mulige neste steg etter forelesningen",
        "bullets": [
            "Velg én notebook og utforsk den i Google Colab",
            "Diskuter et konkret klinisk AI-case i eget fagmiljø",
            "Bruk kursets sjekklister når nye AI-verktøy vurderes lokalt",
        ],
        "references": ["docs/index.html", "docs/for-undervisere.html", "docs/ressurser.html"],
    },
    {
        "type": "content",
        "section": "Avslutning",
        "title": "Takk for oppmerksomheten",
        "bullets": [
            "AI i helse bør ikke bare imponere",
            "Den bør forstås, vurderes og brukes ansvarlig",
            "Det er et klinisk og faglig arbeid, ikke bare et teknologisk prosjekt",
        ],
        "references": ["https://arvidl.github.io/AI-og-helse/", "https://github.com/arvidl/AI-og-helse"],
    },
]


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    for spec in SLIDES:
        slide = prs.slides.add_slide(blank)
        slide_type = spec["type"]

        if slide_type == "title":
            add_title_slide(slide, spec["title"], spec["subtitle"], spec["footer"])
        elif slide_type == "section":
            add_section_slide(slide, spec["title"], spec["subtitle"])
        elif slide_type == "repo_overview":
            add_repo_overview_slide(
                slide,
                spec["title"],
                spec["items"],
                references=spec.get("references"),
                section=spec.get("section"),
            )
        elif slide_type == "case":
            add_case_slide(
                slide,
                spec["title"],
                spec["scenario"],
                spec["left_title"],
                spec["left_points"],
                spec["right_title"],
                spec["right_points"],
                references=spec.get("references"),
                section=spec.get("section"),
            )
        else:
            add_content_slide(
                slide,
                spec["title"],
                spec["bullets"],
                references=spec.get("references"),
                section=spec.get("section"),
            )

    prs.save(OUTPUT_PATH)
    return OUTPUT_PATH


if __name__ == "__main__":
    path = build_presentation()
    print(path)
